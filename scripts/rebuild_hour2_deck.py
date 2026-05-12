"""Add a Copilot Studio Hour 2 block to the agents course deck. NON-DESTRUCTIVE.

Source: docs/_archive/warner-github-agents-REIMAGINED.pptx (81 slides)
Output: warner-github-agents-may-2026.pptx (repo root, 98 slides)

Strategy
- Slide 22 (section divider) is rewritten in-place to preserve the existing visual chrome.
- 17 new Copilot Studio slides are INSERTED between slide 22 and the existing slide 23.
- All 81 original slides are preserved. The existing LangGraph Hour 2 block remains intact
  and continues to flow after the new Copilot Studio block.
- Final ordering: [orig 1..22] + [17 new Copilot Studio slides] + [orig 23..81] = 98 slides.
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SRC = Path("docs/_archive/warner-github-agents-REIMAGINED.pptx")
DST = Path("warner-github-agents-may-2026.pptx")


def set_textbox_text(shape, new_text: str) -> None:
    """Replace a textbox's text, preserving the first run's formatting."""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = new_text
        for r in first_para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first_para.add_run().text = new_text
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)


COPILOT_SLIDES = [
    (
        "HOUR 2 . WHY LOW-CODE FIRST",
        "Why we start in Copilot Studio",
        [
            "Most enterprise agents start in low-code; production rewrites come later",
            "Copilot Studio = opinionated runtime + visual designer + built-in M365 connectors",
            "You can ship to a Teams channel in one afternoon",
            "You will outgrow it, and that is fine. Hour 3 shows what is underneath.",
        ],
        "Set the adoption arc honestly. Most teams ship a Copilot Studio agent before they ship a LangGraph one. That is not a compromise, it is how the funnel works. Then preview Hour 3: same Contoso HR Agent, code version, next hour.",
    ),
    (
        "HOUR 2 . COPILOT STUDIO ANATOMY",
        "The five building blocks",
        [
            "Topics: conversation flows triggered by phrases or events",
            "Knowledge: grounded sources (SharePoint, websites, Dataverse, uploaded files)",
            "Actions: connectors that call out (Outlook, Teams, Power Automate, custom)",
            "Variables: topic-scope and global-scope state",
            "Entities: slot-fill extractors (prebuilt + custom)",
        ],
        "Quick orientation. Do not dwell. Most of the audience has seen Copilot Studio; the build is where the value is. 90 seconds max.",
    ),
    (
        "HOUR 2 . THE AGENT WE ARE BUILDING",
        "Contoso HR Agent: same job, different runtime",
        [
            "Same name and use case as the code version we will see in Hour 3",
            "Knowledge: Microsoft Learn HR-shaped governance content (live, public, no auth)",
            "Topics: Evaluate Resume / Ask HR Policy / Show Last Candidate",
            "Action: Outlook Send email v2 on decision",
            "Output: Adaptive Card with structured evaluation",
            "Publish target: Microsoft Teams",
        ],
        "Set expectations. Show the FINISHED agent first (already-built tenant) so they know what success looks like. Then we walk the build together.",
    ),
    (
        "HOUR 2 . LIVE BUILD . KNOWLEDGE SOURCE",
        "Step 1: add the knowledge source",
        [
            "Add a public-website knowledge source",
            "URL: learn.microsoft.com/training/paths/m365-copilot-administrator/",
            "Why Microsoft Learn over SharePoint: rock-solid, no auth, deterministic, indexed",
            "Verify with one test query in the test pane",
        ],
        "Conference Wi-Fi is hostile. SharePoint introduces tenant scoping and conditional access risk. State explicitly why we picked the public URL. Learners will ask.",
    ),
    (
        "HOUR 2 . LIVE BUILD . GLOBAL VARIABLE",
        "Step 2: create the global variable",
        [
            "Create Global.LastCandidate (String, UseInAIContext on)",
            "Set in Topic 1 step 4 (after slot-fill)",
            "Read in Topic 3 and surfaced to the orchestrator",
            "Demonstrates cross-topic state, the question I get every session",
        ],
        "Pause here. Global vs. topic variables is the #1 question in every Copilot Studio session. Topic-scope dies at end of topic; global-scope lives the whole conversation.",
    ),
    (
        "HOUR 2 . LIVE BUILD . TOPIC 1 (EVALUATE RESUME)",
        "Step 3: build the resume intake topic",
        [
            "Trigger phrases: evaluate this resume / score this candidate / is this candidate a fit",
            "Sequence: 3 Question nodes, SetVariable, SearchAndSummarizeContent, Adaptive Card, Outlook action",
            "Entities: StringPrebuilt x2 + NumberPrebuiltEntity x1",
            "Grounded answer uses the HR Policy Library knowledge source",
        ],
        "Longest build step, about 10 minutes. Narrate each click. Slot-fill rhythm: trigger, ask name, ask role, ask years, set global, generate, render card, email.",
    ),
    (
        "HOUR 2 . LIVE BUILD . ADAPTIVE CARD",
        "Step 4: the structured response",
        [
            "FactSet for candidate / role / years / disposition",
            "TextBlock for reasoning",
            "Action.Submit for Email HR",
            "Cards beat plain text: structured, scannable, accessible, reusable",
        ],
        "Show JSON, then show rendered card side by side. Learners lightbulb moment usually fires here. Reinforce that cards are reusable across Teams, Outlook actionable messages, and Power Apps.",
    ),
    (
        "HOUR 2 . LIVE BUILD . TOPIC 2 (ASK HR POLICY)",
        "Step 5: policy Q&A topic",
        [
            "Trigger phrases: what is our policy on / hr policy / company guideline",
            "Single SearchAndSummarizeContent node, grounded",
            "responseCaptureType: FullResponse for citations",
            "Demonstrates knowledge reuse across topics",
        ],
        "Fastest topic to build, about 4 minutes. Point out you are reusing the same knowledge source from Topic 1. Knowledge is agent-scoped, not topic-scoped.",
    ),
    (
        "HOUR 2 . LIVE BUILD . TOPIC 3 (SHOW LAST CANDIDATE)",
        "Step 6: recall the last candidate",
        [
            "Trigger phrases: who was the last candidate / recap last applicant",
            "ConditionGroup on =IsBlank(Global.LastCandidate)",
            "Empty branch: friendly fallback message",
            "Set branch: surface the global variable",
        ],
        "This topic earns the global-variable slide its keep. Run BEFORE Topic 1 (should say no candidate yet). Then run Topic 1. Then run Topic 3 again. Live state persistence demo.",
    ),
    (
        "HOUR 2 . LIVE BUILD . CONNECTOR ACTION",
        "Step 7: wire the Outlook connector",
        [
            "Office 365 Outlook, Send email v2 (GA, not preview)",
            "Fires at end of Topic 1, after the Adaptive Card renders",
            "Connection authorized in advance, no live OAuth dance",
            "Why not Teams Post message: bot framework channel registration is a stage trap",
        ],
        "Call out pre-authorization explicitly. When learners build this on their own, the first-time OAuth popup can hide behind the browser and stall the topic. Authorize once before you demo.",
    ),
    (
        "HOUR 2 . ORCHESTRATION GOTCHA",
        "The one thing that breaks in week 2",
        [
            "Generative orchestration picks topics by phrase similarity + modelDescription",
            "Overlapping phrases (review this) can route to the wrong topic",
            "Fix: tight modelDescription per topic, semantic not syntactic",
            "Test routing in the test pane BEFORE you publish",
        ],
        "Number-one trip-up. Spend 30 seconds on modelDescription. Show a good one vs. a bad one if time allows. Example: triggers only when the user provides a NEW candidate vs. triggers when the user asks about a PREVIOUSLY evaluated candidate.",
    ),
    (
        "HOUR 2 . PUBLISH TO TEAMS",
        "Step 8: publish and run end-to-end",
        [
            "Channels, Microsoft Teams, Turn on",
            "Approval flow if your tenant requires admin consent",
            "Publish, Get sharing link, Install in a channel",
            "End-to-end: resume eval in Teams, Adaptive Card, email lands in inbox",
        ],
        "The payoff. Have Teams pre-pinned to a second monitor. Run the full flow: trigger Topic 1 in Teams, slot-fill, see the card, switch to Outlook, show the email. About 3 minutes if everything is warm.",
    ),
    (
        "HOUR 2 . DEBUG . WHEN IT DOES NOT WORK",
        "Where to look first",
        [
            "Test pane = your friend. Watch the transcript and variable values live.",
            "Topic did not trigger, check trigger phrases + modelDescription",
            "Generative answer empty, knowledge source not indexed yet (wait, or re-add)",
            "Action failed, check connector run history in Power Automate",
            "Adaptive Card blank, JSON syntax (validate in Adaptive Cards Designer)",
        ],
        "Spend a full 3-4 minutes here. Debug is where learners get stuck on their own. Show the test pane. Show one intentional failure (e.g., comment out the SetVariable, ask Topic 3, get the empty-state message).",
    ),
    (
        "HOUR 2 . LOW-CODE VS. CODE-FIRST",
        "Pick the right tool, both are valid",
        [
            "Stay in Copilot Studio when: M365-first audience, business-user authors, standard connectors, conversational UI, fast iteration",
            "Move to code-first when: custom orchestration, non-M365 ecosystems, custom RAG with embedded ranking, deep observability, CI/CD for the agent itself",
            "Most production landscapes have both",
            "Wrong question: which one. Right question: which one for this agent.",
        ],
        "Bridge to the rest of Hour 2 (LangGraph code version). Both sides are valid. Most production landscapes have both. Do not frame as Copilot Studio vs. LangGraph. Frame as choosing per agent.",
    ),
    (
        "HOUR 2 . WHAT WE BUILT IN COPILOT STUDIO",
        "The Contoso HR Agent, low-code edition",
        [
            "3 topics / 1 connector action / 1 Adaptive Card / 1 global variable",
            "Knowledge-grounded, state-aware",
            "Published to Microsoft Teams",
            "Email on decision",
            "Build time: about 30 minutes",
        ],
        "Quick proud-moment slide. Then transition: Now let us see the SAME agent in code. Coming up: LangGraph + CrewAI version of the Contoso HR Agent.",
    ),
    (
        "HOUR 2 . THE BRIDGE",
        "Same agent, two runtimes",
        [
            "What you just saw: Copilot Studio (no code, M365-native)",
            "What you will see next: LangGraph + CrewAI + ChromaDB + FastMCP",
            "Same domain (HR resume screening), same dispositions, much more control",
            "Why both: production landscapes blend low-code and code-first",
        ],
        "Critical bridge slide. Do not skip. Tell learners: 'You just shipped an agent. Now I am going to show you the same agent built the hard way, and you will see exactly when the hard way is worth it.' Then the existing slide 23 (HOUR 2 . THE APP) takes over.",
    ),
    (
        "",
        "Questions on the low-code build?",
        [
            "Quick check before we open the code",
            "techtrainertim.com   /   @TechTrainerTim",
        ],
        "Optional mini-Q&A before pivoting to the LangGraph block. If running tight on time, skip this slide and go straight into the existing slide 23.",
    ),
]


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"Source deck not found: {SRC}")

    prs = Presentation(str(SRC))
    print(f"Source: {SRC}")
    print(f"Source slide count: {len(prs.slides)}")

    # 1. Rewrite slide 22 (section divider) in place
    slide22 = prs.slides[21]
    for shape in slide22.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt == "Run the Parallel Pipeline":
            set_textbox_text(shape, "Low-Code then Code-First")
        elif txt.startswith("LangGraph + CrewAI"):
            set_textbox_text(
                shape,
                "Copilot Studio Contoso HR Agent, then the same agent in LangGraph + CrewAI.",
            )
    print("Slide 22 (section divider) rewritten in place")

    # 2. Append 17 new Copilot Studio slides (no deletions)
    layout_tc = prs.slide_layouts[1]  # Title and Content
    for kicker, title, bullets, notes in COPILOT_SLIDES:
        slide = prs.slides.add_slide(layout_tc)

        title_ph = slide.shapes.title
        if kicker:
            title_ph.text_frame.text = kicker
            for r in title_ph.text_frame.paragraphs[0].runs:
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0x76, 0x76, 0x76)
            p2 = title_ph.text_frame.add_paragraph()
            p2.text = title
            for r in p2.runs:
                r.font.size = Pt(32)
                r.font.bold = True
        else:
            title_ph.text_frame.text = title
            for r in title_ph.text_frame.paragraphs[0].runs:
                r.font.size = Pt(32)
                r.font.bold = True

        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break
        if body_ph is not None:
            tf = body_ph.text_frame
            tf.text = bullets[0]
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(18)
            for bullet in bullets[1:]:
                para = tf.add_paragraph()
                para.text = bullet
                para.level = 0
                for r in para.runs:
                    r.font.size = Pt(18)

        slide.notes_slide.notes_text_frame.text = notes

    print(f"Appended {len(COPILOT_SLIDES)} Copilot Studio slides (originals preserved)")
    print(f"Slide count after append: {len(prs.slides)}")

    # 3. Reorder: insert the new block immediately after slide 22 (section divider)
    slide_id_lst = prs.slides._sldIdLst
    all_ids = list(slide_id_lst)
    pre = all_ids[:22]              # slides 1..22 (Hour 1 + Hour 2 section divider)
    original_rest = all_ids[22:81]  # original slides 23..81 (Hour 2 LangGraph through end)
    new_block = all_ids[81:]        # the 17 new slides we just appended

    print(f"Reorder counts: pre={len(pre)} new={len(new_block)} original_rest={len(original_rest)}")
    assert len(new_block) == len(COPILOT_SLIDES), "new_block size mismatch"

    for s in all_ids:
        slide_id_lst.remove(s)
    for s in pre + new_block + original_rest:
        slide_id_lst.append(s)

    # 4. Save to repo root
    prs.save(str(DST))
    size_kb = os.path.getsize(DST) // 1024
    print(f"\nWritten: {DST} ({size_kb} KB)")

    # 5. Verify
    verify = Presentation(str(DST))
    print(f"\nVerification: {len(verify.slides)} slides total (expected 98)")
    print("\nFull slide list:")
    for i, s in enumerate(verify.slides, 1):
        title = ""
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0][:75]
                break
        print(f"  {i:3}. {title}")


if __name__ == "__main__":
    main()
