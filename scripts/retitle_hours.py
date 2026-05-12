"""Retitle kicker text on `warner-github-agents-may-2026.pptx` to match the new structure.

New structure:
  Hour 1 = Agents concepts          (slides 10-21, kickers already "HOUR 1 ·", no change)
  Hour 2 = Copilot Studio           (slides 23-39, new block, kickers already "HOUR 2 ·", no change)
  Hour 3 = LangGraph code walkthrough  (slides 40-57, currently "HOUR 2 ·" -> retitle to "HOUR 3 ·")
  Hour 4 = MCP + deployment         (slides 59-97, currently mix of "HOUR 3 ·" / "HOUR 4 ·"
                                     -> retitle ALL to "HOUR 4 ·")

Section dividers:
  Slide 40 group: the "HOUR 2 - what we covered" recap at slide 57 should become "HOUR 3 - what we covered"
  Slide 58 is the "03" divider -> change to "04 reserved" no, slide 58 is the existing "03" divider for the
    old MCP-led Hour 3. Under the new plan, slide 58 becomes the "HOUR 3 -> HOUR 4 transition" cap.
  Slide 73 ("Hour 3 - what we covered") becomes mid-Hour-4 recap.
  Slide 74 is "04" divider for the production half -> rename divider subtitle to reflect MCP -> deployment flow.

This script touches ONLY text content of existing shapes. No deletions, no inserts.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

DECK = Path("warner-github-agents-may-2026.pptx")


def set_textbox_text(shape, new_text: str) -> None:
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = new_text
        for r in first_para.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first_para.add_run().text = new_text
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)


def rewrite_text(shape, old_text: str, new_text: str) -> bool:
    """If the shape's exact stripped text matches old_text, replace it."""
    if not shape.has_text_frame:
        return False
    if shape.text_frame.text.strip() == old_text.strip():
        set_textbox_text(shape, new_text)
        return True
    return False


def replace_prefix(shape, old_prefix: str, new_prefix: str) -> bool:
    """If the shape's text starts with old_prefix, swap the prefix."""
    if not shape.has_text_frame:
        return False
    txt = shape.text_frame.text
    if txt.startswith(old_prefix):
        new_text = new_prefix + txt[len(old_prefix):]
        set_textbox_text(shape, new_text)
        return True
    return False


def main() -> None:
    prs = Presentation(str(DECK))
    total = len(prs.slides)
    print(f"Deck: {DECK}  ({total} slides)")

    changes = 0

    # --- Hour 3: slides 40-57 are the LangGraph block (currently labeled HOUR 2) ---
    # These were originally Hour 2 LangGraph slides; under the new plan they are Hour 3.
    for slide_idx in range(40, 58):  # 1-indexed inclusive 40..57
        slide = prs.slides[slide_idx - 1]
        for shape in slide.shapes:
            if replace_prefix(shape, "HOUR 2 ", "HOUR 3 "):
                changes += 1
            # Recap slide 57 has "Hour 2 - what we covered"
            elif rewrite_text(shape, "Hour 2 — what we covered", "Hour 3 — what we covered"):
                changes += 1

    # --- Slide 58: section divider currently labeled "03" with subtitle about MCP ---
    # New plan: this is the transition INTO Hour 4 (MCP + deployment).
    slide58 = prs.slides[57]
    for shape in slide58.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt == "03":
            set_textbox_text(shape, "04")
            changes += 1
        elif txt == "Tools, Knowledge & MCP":
            set_textbox_text(shape, "MCP and Production Deployment")
            changes += 1
        elif txt.startswith("FastMCP 2"):
            set_textbox_text(shape, "Contoso HR Agent + MCP Inspector, then deploy to Azure with guardrails and observability.")
            changes += 1

    # --- Hour 4: slides 59-72 are the MCP block (currently labeled HOUR 3) ---
    for slide_idx in range(59, 74):
        slide = prs.slides[slide_idx - 1]
        for shape in slide.shapes:
            if replace_prefix(shape, "HOUR 3 ", "HOUR 4 "):
                changes += 1
            elif rewrite_text(shape, "Hour 3 — what we covered", "Hour 4 — MCP recap"):
                changes += 1

    # --- Slide 74: was the "04 - Production on Azure" divider ---
    # Under the new plan, Hour 4 starts at slide 58. Slide 74 is now an in-hour pivot card
    # from MCP to deployment. Replace the "04" with "Pivot" and rewrite the subtitle.
    slide74 = prs.slides[73]
    for shape in slide74.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt == "04":
            set_textbox_text(shape, "Pivot")
            changes += 1
        elif txt == "Production on Azure":
            set_textbox_text(shape, "From MCP to production")
            changes += 1
        elif txt.startswith("Foundry") or "APIM" in txt or "Container Apps" in txt:
            set_textbox_text(
                shape,
                "MCP demo is done. Now deploy the same Contoso HR Agent to Azure with guardrails, observability, and cost controls.",
            )
            changes += 1

    # --- Hour 4 (production half): slides 75-95 ---
    # These already say HOUR 4. No change needed for kickers. But slide 93 says
    # "What you built today" which is fine. Skip.

    # Save in place
    prs.save(str(DECK))
    print(f"\nTotal kicker rewrites: {changes}")
    print(f"Saved: {DECK}")

    # Verify
    v = Presentation(str(DECK))
    print(f"\nVerification: {len(v.slides)} slides")
    print("\nKicker check (slides 39-75):")
    for i, s in enumerate(v.slides, 1):
        if 39 <= i <= 75:
            title = ""
            for shape in s.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().split("\n")[0][:75]
                    break
            print(f"  {i:3}. {title}")


if __name__ == "__main__":
    main()
